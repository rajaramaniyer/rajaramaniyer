#!/opt/homebrew/opt/python@3.9/libexec/bin/python

import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import sys
import subprocess
import os
import glob
import re
import os.path
from os import path

def get_lines(filename):
  file = open(filename, 'r', encoding="UTF-8")
  Lines = file.readlines()
  file.close()
  return Lines

def remspace(my_str):
  if len(my_str) < 2: # returns ' ' unchanged
    return my_str
  if my_str[-1] == '\n':
    if my_str[-2] == ' ':
      return my_str[:-2] + '\n'
  if my_str[-1] == ' ':
    return my_str[:-1]
  return my_str

def scrub_file(filename):
  Lines = get_lines(filename)
  # Scrub Lines
  file = open(filename, "wt", encoding="UTF-8")
  for line in Lines:
      line=line.replace("||", "॥").replace("।।", "॥").replace("|","।")
      line=line.replace("॥"," ॥ ").replace("।"," । ")
      line=re.sub(r"\t"," ",line)
      line=line.replace(".",". ")
      line=line.replace("-"," - ")
      line=line.replace(" "," ")
      while ( re.search("  ", line) ):
        line=re.sub("  "," ",line)
      line=remspace(line)
      file.write(line)
  file.close()

#
# Main Code Starts
#
PPTFILE="generated-ppt.pptx"
sanskrit_file="/Users/rajaramaniyer/books/srisrianna/dwadasha-nama-keerthanams-sanskrit.txt"

if not path.exists(sanskrit_file):
  input("Sanskrit file does not exists. " + sanskrit_file)
  sys.exit()

scrub_file(sanskrit_file)

sanskrit_lines=get_lines(sanskrit_file)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

text_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(text_slide_layout)

for shape in slide.placeholders:
  print('%d %s' % (shape.placeholder_format.idx, shape.name))


slide.placeholders[0].width = Inches(1)
slide.placeholders[0].text = 'Slide 1'
slide.placeholders[1].text = 'Page 1'

prs.save(PPTFILE)
print("run open %s" % PPTFILE)