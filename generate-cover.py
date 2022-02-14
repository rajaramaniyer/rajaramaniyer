import collections.abc
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import sys
import csv
import subprocess
import os
import glob
import time
import os.path
from os import path

#
# Main Code Starts
#API_KEY="AIzaSyA5VQp8hffcJKWDxDmP7mST7pbb78fSzpc"
#
prs = Presentation('template.pptm')
title_slide = prs.slide_layouts[0]

files = glob.glob('C:/Users/rajar/Pictures/Bhagavatham10/??.png')
for f in files:
    adyayam_number=os.path.splitext(os.path.basename(f))[0]
    titletext='Srimad Bhagavatham | Canto 10 | Chapter ' + adyayam_number

    picture_file_name = 'C:/Users/rajar/Pictures/Bhagavatham10/' + adyayam_number + '.png'
    if not path.exists(picture_file_name):
        input("Picture file does not exists. " + picture_file_name)
        sys.exit()

    slide = prs.slides.add_slide(title_slide)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = titletext
    picture = shapes.placeholders[10].insert_picture(picture_file_name)
    shapes.placeholders[15].text_frame.text = "श्रीमद्भागवतमहापुराणम्"
    shapes.placeholders[14].text_frame.text = "दशमः स्कन्दः"
    body_shape = shapes.placeholders[12] #Adyayam Text added later
    shapes.placeholders[16].text_frame.text = "10"
    shapes.placeholders[17].text_frame.text = adyayam_number

    ## Main Loop to add content slides
    lyrics_file='C:/Users/rajar/srimad_bhaghavatham/sanskrit/canto10/chapter'+adyayam_number+".txt"
    if not path.exists(lyrics_file):
        input("Picture file does not exists. " + picture_file_name)
        sys.exit()

    file = open(lyrics_file, 'r', encoding="UTF-8")
    Lines = file.readlines()
    file.close()
    body_shape.text_frame.text = Lines[0]

prs.save('generated-ppt.pptm')

next_action=input ("Export JPG (e) / Open (o): ")
if next_action == "e":
  files = glob.glob('Slide????.jpg')
  for f in files:
    try:
        os.remove(f)
    except OSError as e:
        print("Error: %s : %s" % (f, e.strerror))
  p = subprocess.Popen(["C:/Program Files/Microsoft Office 15/root/office15/POWERPNT.EXE", "/M", "generated-ppt.pptm", "Save_PowerPoint_Slide_as_Images"])
elif next_action == "o":
  p = subprocess.Popen(["C:/Program Files/Microsoft Office 15/root/office15/POWERPNT.EXE", "generated-ppt.pptm"])
