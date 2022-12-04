import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import sys
import subprocess
import os
import glob
import re
import os.path
from os import path

def get_lines(filename):
  file = open(filename, 'r', encoding="UTF-8")
  Lines = file.readlines()
  file.close()
  return Lines

def remspace(my_str):
  if len(my_str) < 2: # returns ' ' unchanged
    return my_str
  if my_str[-1] == '\n':
    if my_str[-2] == ' ':
      return my_str[:-2] + '\n'
  if my_str[-1] == ' ':
    return my_str[:-1]
  return my_str

def scrub_file(filename):
  Lines = get_lines(filename)
  # Scrub Lines
  file = open(filename, "wt", encoding="UTF-8")
  for line in Lines:
      line=line.replace("||", "॥").replace("।।", "॥").replace("|","।")
      line=line.replace("॥"," ॥ ").replace("।"," । ")
      line=re.sub(r"\t"," ",line)
      line=line.replace(".",". ")
      line=line.replace("-"," - ")
      line=line.replace(" "," ")
      while ( re.search("  ", line) ):
        line=re.sub("  "," ",line)
      line=remspace(line)
      file.write(line)
  file.close()

#
# Main Code Starts
#
if len(sys.argv) != 2:
  song_number=int(input("Song Number: "))
else:
  song_number=int(sys.argv[1])
SONG_PREFIX="song" + str(song_number).zfill(2)
USERPROFILE=os.environ['USERPROFILE']
TDFOLDER=USERPROFILE + "\\books\\srisrianna\\tamil-divyanamam\\"
PPTFILE=USERPROFILE + "\\Videos\\tamil-divyanamam\\" + SONG_PREFIX + ".pptm"
TEMPLATEFILE=USERPROFILE + "\\Videos\\tamil-divyanamam\\template.pptm"
tamil_file=TDFOLDER + SONG_PREFIX + ".txt"

if not path.exists(tamil_file):
  input("Tamil file does not exists. " + tamil_file)
  sys.exit()

if not path.exists(TEMPLATEFILE):
  input("Powerpoint Template file does not exists. " + TEMPLATEFILE)
  sys.exit()

scrub_file(tamil_file)

tamil_lines=get_lines(tamil_file)

prs = Presentation(TEMPLATEFILE)
title_slide = prs.slide_layouts[0]

##
## For checking contents of the slide
#slide = prs.slides.add_slide(title_slide)
#for shape in slide.placeholders:
#  print('%d %s' % (shape.placeholder_format.idx, shape.name))

for t in tamil_lines:
  t_trimmed=re.sub("\n","",t)
  if len(t_trimmed) > 0:
    slide = prs.slides.add_slide(title_slide)
    shapes = slide.shapes
    tamil_p = shapes.placeholders[10].text_frame
    tamil_text = tamil_p.paragraphs[0]
    tamil_text.text = t_trimmed

prs.save(PPTFILE)

files = glob.glob(SONG_PREFIX + "-Slide????.jpg")
for f in files:
  try:
      os.remove(f)
  except OSError as e:
      print("Error: %s : %s" % (f, e.strerror))
p = subprocess.Popen(["C:/Program Files/Microsoft Office 15/root/office15/POWERPNT.EXE", "/M", PPTFILE, "TDSave_PowerPoint_Slide_as_Images"])
