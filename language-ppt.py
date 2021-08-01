from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import sys
import csv
import subprocess
import os
import glob
import time

#
# Main Code Starts
#
sanskritcsv = csv.reader(open('sanskrit.csv', encoding="UTF-8"))
tamilcsv = csv.reader(open('tamil.csv', encoding="UTF-8"))
prs = Presentation('template.pptm')
template = [prs.slide_layouts[1],prs.slide_layouts[2],prs.slide_layouts[3]]

#
# For checking contents of the slide
#for shape in slide.placeholders:
#  print('%d %s' % (shape.placeholder_format.idx, shape.name))

lines=0
slidecount=len(prs.slides)
iTemplate=0
for s in sanskritcsv:
  t = tamilcsv.__next__()
  if lines == 0:
    if len(s) > 0 and "श्रीगुरोरष्टकम्‌" == s[0]:
      iTemplate = 2
    if iTemplate > 0:
      if iTemplate == 1:
        iTemplate = 2
      else:
        iTemplate = 1
    slide = prs.slides.add_slide(template[iTemplate])
    slidecount = slidecount+1
    shapes = slide.shapes
    placeholders = shapes.placeholders
    sanskrit_p = placeholders[13].text_frame.paragraphs[0]
    tamil_p = placeholders[1].text_frame.paragraphs[0]
  sanskrit_text = sanskrit_p.add_run()
  tamil_text = tamil_p.add_run()
  if len(s) > 0:
    sanskrit_text.text = s[0]
  if len(t) > 0:
    tamil_text.text = t[0]
    if len(s) == 0:
      font = tamil_text.font
      font.size = Pt(18)
      font.color.rgb = RGBColor(0x00, 0x00, 0xFF)

  if sanskrit_text.text.find("-") != -1 and len(sanskrit_text.text) - sanskrit_text.text.find("-") > 1:
    font = sanskrit_text.font
    font.color.rgb = RGBColor(0x00, 0x00, 0xFF)
    font = tamil_text.font
    font.color.rgb = RGBColor(0x00, 0x00, 0xFF)

  sanskrit_text.text = sanskrit_text.text + "\n"
  tamil_text.text = tamil_text.text + "\n"
  lines = lines + 1
  if lines > 6:
    lines = 0

filename="Poojai.pptm"
prs.save(filename)

print ("Generated. Slide Count = ", slidecount)
next_action=input ("Export JPG (e) / Open (o): ")
if next_action == "e":
  files = glob.glob('Slide????.???')
  for f in files:
    try:
        os.remove(f)
    except OSError as e:
        print("Error: %s : %s" % (f, e.strerror))
  p = subprocess.Popen(["C:/Program Files/Microsoft Office 15/root/office15/POWERPNT.EXE", "/M", filename, "Save_PowerPoint_Slide_as_Images"])
  while not next(glob.iglob("Slide*" + str(slidecount) + ".jpg"), None):
    time.sleep(1)
  p.terminate()
  p.wait()
elif next_action == "o":
  p = subprocess.Popen(["C:/Program Files/Microsoft Office 15/root/office15/POWERPNT.EXE", filename])

# gci Slide*.jpg | ren -n {[regex]::replace($_.Name, '\d+', {"$args".PadLeft(4, '0')})}
# भोगः
# பூஜை