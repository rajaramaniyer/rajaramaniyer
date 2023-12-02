import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import sys
import subprocess
import os
import glob
import re
import os.path
from os import path
import wave, contextlib
import uuid
import time, math
import xml.etree.ElementTree as ET
from xml.dom.minidom import parseString

def get_lines(filename):
  file = open(filename, 'r', encoding="UTF-8")
  Lines = file.readlines()
  file.close()
  return Lines

def remspace(my_str):
  if len(my_str) < 2: # returns ' ' unchanged
    return my_str
  if my_str[-1] == '\n':
    if my_str[-2] == ' ':
      return my_str[:-2] + '\n'
  if my_str[-1] == ' ':
    return my_str[:-1]
  return my_str

def scrub_file(filename):
  Lines = get_lines(filename)
  # Scrub Lines
  file = open(filename, "wt", encoding="UTF-8")
  for line in Lines:
      line=line.replace("||", "॥").replace("।।", "॥").replace("|","।")
      line=line.replace("॥"," ॥ ").replace("।"," । ")
      line=re.sub(r"\t"," ",line)
      line=line.replace(".",". ")
      line=line.replace("-"," - ")
      line=line.replace(" "," ")
      while ( re.search("  ", line) ):
        line=re.sub("  "," ",line)
      line=remspace(line)
      file.write(line)
  file.close()

def build_shathakam_ppt(short_name='', long_name='premikendra-sadguru-stuti'):
  song_number = get_song_number(long_name)
  SONG_PREFIX=song_number
  USERPROFILE=os.environ['USERPROFILE']
  RSFOLDER=USERPROFILE + "\\books\\srisrianna\\"+long_name+"\\"
  IMGFILE=USERPROFILE + "\\Pictures\\001 "+long_name+"\\YouTube\\" + SONG_PREFIX + ".*"
  SLIDEIMGFILE=USERPROFILE + "\\Pictures\\001 "+long_name+"\\YouTube\\" + SONG_PREFIX
  IMGFILES=glob.glob(IMGFILE)
  PPTFILE=USERPROFILE + "\\Videos\\"+long_name+"\\" + SONG_PREFIX + ".pptm"
  TEMPLATEFILE=USERPROFILE + "\\Videos\\"+long_name+"\\template.pptm"
  sanskrit_file=RSFOLDER + SONG_PREFIX + "-sanskrit.txt"
  tamil_file=RSFOLDER + SONG_PREFIX + "-tamil.txt"

  if not path.exists(sanskrit_file):
    input("Sanskrit file does not exists. " + sanskrit_file)
    sys.exit()

  if not path.exists(tamil_file):
    input("Tamil file does not exists. " + tamil_file)
    sys.exit()

  if not path.exists(TEMPLATEFILE):
    input("Powerpoint Template file does not exists. " + TEMPLATEFILE)
    sys.exit()

  if len(IMGFILES) < 1:
    input("Image file does not exists. " + IMGFILE)
    sys.exit()

  if len(IMGFILES) > 1:
    input("More than 1 Image file found: ")
    for img in IMGFILES:
      print(img)
    sys.exit()

  scrub_file(sanskrit_file)
  scrub_file(tamil_file)

  sanskrit_lines=get_lines(sanskrit_file)
  tamil_lines=get_lines(tamil_file)

  if len(sanskrit_lines) != len(tamil_lines):
    input("Sanskrit file has " + str(len(sanskrit_lines)) + " lines, while Tamil file has " + str(len(tamil_lines)) + ". they are not having same count of lines.")
    sys.exit()

  prs = Presentation(TEMPLATEFILE)
  text_slide = prs.slide_layouts[1]

  ##
  ## For checking contents of the slide
  #slide = prs.slides.add_slide(text_slide)
  #for shape in slide.placeholders:
  #  print('%d %s' % (shape.placeholder_format.idx, shape.name))

  titleline=-1
  line=0
  count=0
  slide_count=0
  for s in sanskrit_lines:
    index=s.find("-")
    if index != -1 and len(s) - index > 2 and index > 2 :
      titleline = line
    else:
      titleline = -1
    t = tamil_lines[line]
    if count == 0:
      slide_count += 1
      slide = prs.slides.add_slide(text_slide)
      shapes = slide.shapes
      SLIDEIMGFILES=glob.glob(SLIDEIMGFILE + "-" + str(slide_count).zfill(4) + ".*")
      try:
        if len(SLIDEIMGFILES) == 0:
          SLIDEIMGFILES = IMGFILES
        picture = shapes.placeholders[10].insert_picture(SLIDEIMGFILES[0])
      except ValueError as e:
        print("\n\nError adding image '%s'\n\n" % SLIDEIMGFILES[0])
        print(e)
        input("\nPress any key to exit\n\n")
        sys.exit()
      sanskrit_p = shapes.placeholders[11].text_frame
      tamil_p = shapes.placeholders[12].text_frame
      sanskrit_text = sanskrit_p.paragraphs[0]
      tamil_text = tamil_p.paragraphs[0]
    else:
      sanskrit_text = sanskrit_p.add_paragraph()
      tamil_text = tamil_p.add_paragraph()
    if len(s) > 0:
      sanskrit_text.text = re.sub("\n","",s)
      if s.lstrip().find("-") == 0:
        sanskrit_text.alignment = PP_ALIGN.RIGHT
    if len(t) > 0:
      tamil_text.text = re.sub("\n","",t)
      if t.lstrip().find("-") == 0:
        tamil_text.alignment = PP_ALIGN.RIGHT
    if titleline != -1:
      font = sanskrit_text.font
      font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)
      font = tamil_text.font
      font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)
    count = count + 1
    if count > 4:
      count = 0
    line = line + 1

  prs.save(PPTFILE)

  files = glob.glob(SONG_PREFIX + "-Slide????.jpg")
  for f in files:
    try:
        os.remove(f)
    except OSError as e:
        print("Error: %s : %s" % (f, e.strerror))
  # p = subprocess.Popen(["C:/Program Files/Microsoft Office 15/root/office15/POWERPNT.EXE", "/M", PPTFILE, "PSSSave_PowerPoint_Slide_as_Images"])
  p = subprocess.Popen(["C:/Program Files/LibreOffice/program/soffice.exe", "--headless", PPTFILE, "macro:///Standard.Module1.WritePNGSlides()"])
  p.wait()

def get_song_number(long_name):
  if len(sys.argv) != 2:
    song_number=input("%s Song Number: " % long_name)
  else:
    song_number=sys.argv[1]
  return song_number

build_shathakam_ppt()