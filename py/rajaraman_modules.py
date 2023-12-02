import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import sys
import subprocess
import os
import glob
import re
import os.path
from os import path
import wave, contextlib
import uuid
import time, math
import xml.etree.ElementTree as ET
from xml.dom.minidom import parseString
import http.client
import httplib2
import requests
import random
import os
import time
import pickle
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
import os
import sys
from os import path

def get_lines(filename):
  file = open(filename, 'r', encoding="UTF-8")
  Lines = file.readlines()
  file.close()
  return Lines

def remspace(my_str):
  if len(my_str) < 2: # returns ' ' unchanged
    return my_str
  if my_str[-1] == '\n':
    if my_str[-2] == ' ':
      return my_str[:-2] + '\n'
  if my_str[-1] == ' ':
    return my_str[:-1]
  return my_str

def scrub_file(filename):
  Lines = get_lines(filename)
  # Scrub Lines
  file = open(filename, "wt", encoding="UTF-8")
  for line in Lines:
      line=line.replace("||", "॥").replace("।।", "॥").replace("|","।")
      line=line.replace("॥"," ॥ ").replace("।"," । ")
      line=re.sub(r"\t"," ",line)
      line=line.replace(".",". ")
      line=line.replace("-"," - ")
      line=line.replace(" "," ")
      while ( re.search("  ", line) ):
        line=re.sub("  "," ",line)
      line=remspace(line)
      file.write(line)
  file.close()

def build_shathakam_ppt(short_name='GS', long_name='govindashathakam'):
  song_number = get_song_number(long_name)
  SONG_PREFIX=short_name + str(song_number).zfill(3)
  USERPROFILE=os.environ['USERPROFILE']
  RSFOLDER=USERPROFILE + "\\books\\srisrianna\\"+long_name+"\\"
  IMGFILE=USERPROFILE + "\\Pictures\\001 "+long_name+"\\YouTube\\" + SONG_PREFIX + ".*"
  SLIDEIMGFILE=USERPROFILE + "\\Pictures\\001 "+long_name+"\\YouTube\\" + SONG_PREFIX
  IMGFILES=glob.glob(IMGFILE)
  PPTFILE=USERPROFILE + "\\Videos\\"+long_name+"\\" + SONG_PREFIX + ".pptm"
  TEMPLATEFILE=USERPROFILE + "\\Videos\\"+long_name+"\\template.pptm"
  sanskrit_file=RSFOLDER + SONG_PREFIX + "-sanskrit.txt"
  tamil_file=RSFOLDER + SONG_PREFIX + "-tamil.txt"
  index_file=RSFOLDER + "index.txt"

  if not path.exists(sanskrit_file):
    input("Sanskrit file does not exists. " + sanskrit_file)
    sys.exit()

  if not path.exists(tamil_file):
    input("Tamil file does not exists. " + tamil_file)
    sys.exit()

  if not path.exists(TEMPLATEFILE):
    input("Powerpoint Template file does not exists. " + TEMPLATEFILE)
    sys.exit()

  if len(IMGFILES) < 1:
    input("Image file does not exists. " + IMGFILE)
    sys.exit()

  if len(IMGFILES) > 1:
    input("More than 1 Image file found: ")
    for img in IMGFILES:
      print(img)
    sys.exit()

  index_lines=get_lines(index_file)
  if len(index_lines) < (song_number):
    input("Index file does not have title for song " + str(song_number))
    sys.exit()
  title_text=index_lines[song_number-1].split("|")

  scrub_file(sanskrit_file)
  scrub_file(tamil_file)

  sanskrit_lines=get_lines(sanskrit_file)
  tamil_lines=get_lines(tamil_file)

  if len(sanskrit_lines) != len(tamil_lines):
    input("Sanskrit file has " + str(len(sanskrit_lines)) + " lines, while Tamil file has " + str(len(tamil_lines)) + ". they are not having same count of lines.")
    sys.exit()

  prs = Presentation(TEMPLATEFILE)
  text_slide = prs.slide_layouts[1]
  title_slide = prs.slide_layouts[0]

  ##
  ## For checking contents of the slide
  #slide = prs.slides.add_slide(text_slide)
  #for shape in slide.placeholders:
  #  print('%d %s' % (shape.placeholder_format.idx, shape.name))

  ## First Slide
  slide = prs.slides.add_slide(title_slide)
  shapes = slide.shapes
  shapes.placeholders[0].text_frame.text = str(song_number)
  shapes.placeholders[1].text_frame.text = re.sub("\n","",title_text[0])
  shapes.placeholders[11].text_frame.text = re.sub("\n","",title_text[1])
  picture = shapes.placeholders[10].insert_picture(IMGFILES[0])

  titleline=-1
  line=0
  count=0
  slide_count=1
  for s in sanskrit_lines:
    index=s.find("-")
    if index != -1 and len(s) - index > 2 and index > 2 :
      titleline=line
    if line != titleline:
      t = tamil_lines[line]
      if count == 0:
        slide_count += 1
        slide = prs.slides.add_slide(text_slide)
        shapes = slide.shapes
        SLIDEIMGFILES=glob.glob(SLIDEIMGFILE + "-" + str(slide_count).zfill(4) + ".*")
        try:
          if len(SLIDEIMGFILES) == 0:
            SLIDEIMGFILES = IMGFILES
          picture = shapes.placeholders[10].insert_picture(SLIDEIMGFILES[0])
        except ValueError as e:
          print("\n\nError adding image '%s'\n\n" % SLIDEIMGFILES[0])
          print(e)
          input("\nPress any key to exit\n\n")
          sys.exit()
        if titleline != -1:
          shapes.placeholders[13].text_frame.text = re.sub("\n","",sanskrit_lines[titleline])
          shapes.placeholders[14].text_frame.text = re.sub("\n","",tamil_lines[titleline])
        sanskrit_p = shapes.placeholders[11].text_frame
        tamil_p = shapes.placeholders[12].text_frame
        sanskrit_text = sanskrit_p.paragraphs[0]
        tamil_text = tamil_p.paragraphs[0]
      else:
        sanskrit_text = sanskrit_p.add_paragraph()
        tamil_text = tamil_p.add_paragraph()
      if len(s) > 0:
        sanskrit_text.text = re.sub("\n","",s)
        if s.lstrip().find("-") == 0:
          sanskrit_text.alignment = PP_ALIGN.RIGHT
      if len(t) > 0:
        tamil_text.text = re.sub("\n","",t)
        if t.lstrip().find("-") == 0:
          tamil_text.alignment = PP_ALIGN.RIGHT
      count = count + 1
      if count > 4:
        count = 0
    line = line + 1

  prs.save(PPTFILE)

  files = glob.glob(SONG_PREFIX + "-Slide????.jpg")
  for f in files:
    try:
        os.remove(f)
    except OSError as e:
        print("Error: %s : %s" % (f, e.strerror))
  # subprocess.Popen(["C:/Program Files/Microsoft Office 15/root/office15/POWERPNT.EXE", "/M", PPTFILE, short_name + "Save_PowerPoint_Slide_as_Images"])
  p = subprocess.Popen(["C:/Program Files/LibreOffice/program/soffice.exe", "--headless", PPTFILE, "macro:///Standard.Module1.WritePNGSlides()"])
  p.wait()

def get_song_number(long_name):
  if len(sys.argv) != 2:
    song_number=int(input("%s Song Number: " % long_name))
  else:
    song_number=int(sys.argv[1])
  return song_number

def getTimeString(totalmilliseconds):
    milliseconds=str(math.floor(totalmilliseconds)%1000).zfill(3)
    hours=str(math.floor(totalmilliseconds/1000/60/60)).zfill(2)
    minutes=str(math.floor((totalmilliseconds/1000/60)%60)).zfill(2)
    seconds=str(math.floor((totalmilliseconds/1000)%60)).zfill(2)
    return "{hours}:{minutes}:{seconds}.{milliseconds}".format(hours=hours, minutes=minutes, seconds=seconds, milliseconds=milliseconds)

def build_shathakam_mlt(short_name='GS', long_name='govindashathakam',wavefilefolder=None):
    song_number = get_song_number(long_name)
    SONG_PREFIX=short_name + str(song_number).zfill(3)
    USERPROFILE=os.environ['USERPROFILE']
    if wavefilefolder == None:
       wavefilefolder = USERPROFILE + "/lmms/projects/"
    wavefilename=wavefilefolder + SONG_PREFIX + ".wav"
    if not os.path.exists(wavefilename):
        input(wavefilename + " does not exists")
        sys.exit()

    duration="00:00:00.000"
    durationSeconds=0.0
    with contextlib.closing(wave.open(wavefilename,'r')) as f:
        frames = f.getnframes()
        rate = f.getframerate()
        durationMilliSeconds = round(float(frames) / float(rate) * 1000.0)
        duration = getTimeString(durationMilliSeconds)

    labelTrackFile=USERPROFILE + "\\Videos\\"+long_name+"\\" + SONG_PREFIX + ".txt"
    if not os.path.exists(labelTrackFile):
        input(labelTrackFile + " does not exists")
        sys.exit()

    file1 = open(labelTrackFile, 'r')
    labelTrack = file1.readlines()
    file1.close()

    audioHash=uuid.uuid4().hex

    block1=[]
    block2=[]
    block3=[]
    block4=[]
    block2.append('  <playlist id="main_bin">\n')
    block2.append('    <property name="xml_retain">1</property>\n')
    block2.append('    <entry producer="chain0" in="00:00:00.000" out="{}"/>\n'.format(duration))
    #
    block4.append('  <playlist id="playlist0">\n')
    block4.append('    <property name="shotcut:video">1</property>\n')
    block4.append('    <property name="shotcut:name">V1</property>\n')

    id=0
    files=sorted(glob.glob(USERPROFILE + "\\Videos\\"+long_name+"\\" + SONG_PREFIX+"-Slide????.*"))

    equalslideduration=getTimeString(round(durationMilliSeconds/len(files)))
    if len(labelTrack) > 0 and len(files)-1 != len(labelTrack):
        if len(labelTrack) < len(files):
            prev=-1
            for label in labelTrack:
                durationParts = label.split("\t")
                if prev != -1 and prev != durationParts[0]:
                    print("Line " + str(id+1) + " is having a break")
                    break
                prev=durationParts[1]
                id+=1
        print("WARNING: labelTrack and slideCount are Not in sync.")
        print("There are %d labels while there are %d slides" %(len(labelTrack), len(files)))
        #print("Setting all slide's duration equal, which is %s" %(equalslideduration))
        input ("Press enter to continue: ")

    id=0
    currentTime=time.strftime("%Y-%m-%dT%H:%M:%S",time.localtime())
    prevmarker=0
    prevlabelpos=0.0
    for slidename in files:
        if id < len(labelTrack):
            durationParts = labelTrack[id].split("\t")
            mins=0
            if float(durationParts[1]) - float(durationParts[0]) == 0:
                seconds=float(durationParts[0])-prevlabelpos
                prevlabelpos=float(durationParts[0])
            else:
                seconds=float(durationParts[1]) - float(durationParts[0])
            roundSeconds=math.floor(seconds)
            microframes=(math.floor((seconds-roundSeconds)*1000/25)*25)-25
            if microframes < 0:
                microframes = 600 + microframes
                roundSeconds -= 1
            slideduration="00:%02.0f:%02.0f.%03.0f" %(mins, roundSeconds, microframes)
        else:
            slideduration=equalslideduration
        #
        hashvalue=uuid.uuid4().hex
        #
        block1.append('  <producer id="producer{id}" in="00:00:00.000" out="03:59:59.960">\n'.format(id=id))
        block1.append('    <property name="length">04:00:00.000</property>\n')
        block1.append('    <property name="eof">pause</property>\n')
        block1.append('    <property name="resource">{slidename}</property>\n'.format(slidename=slidename))
        block1.append('    <property name="ttl">1</property>\n')
        block1.append('    <property name="aspect_ratio">1</property>\n')
        block1.append('    <property name="progressive">1</property>\n')
        block1.append('    <property name="seekable">1</property>\n')
        block1.append('    <property name="mlt_service">qimage</property>\n')
        block1.append('    <property name="creation_time">{currentTime}</property>\n'.format(currentTime=currentTime))
        block1.append('    <property name="shotcut:hash">{hashvalue}</property>\n'.format(hashvalue=hashvalue))
        block1.append('    <property name="xml">was here</property>\n')
        block1.append('  </producer>\n')
        #
        block2.append('    <entry producer="producer{id}" in="00:00:00.000" out="00:00:03.960"/>\n'.format(id=id))
        #
        block3.append('  <producer id="producer{id}" in="00:00:00.000" out="03:59:59.960">\n'.format(id=id+len(files)))
        block3.append('    <property name="length">04:00:00.000</property>\n')
        block3.append('    <property name="eof">pause</property>\n')
        block3.append('    <property name="resource">{slidename}</property>\n'.format(slidename=slidename))
        block3.append('    <property name="ttl">1</property>\n')
        block3.append('    <property name="aspect_ratio">1</property>\n')
        block3.append('    <property name="progressive">1</property>\n')
        block3.append('    <property name="seekable">1</property>\n')
        block3.append('    <property name="mlt_service">qimage</property>\n')
        block3.append('    <property name="shotcut:hash">{hashvalue}</property>\n'.format(hashvalue=hashvalue))
        block3.append('    <property name="creation_time">2021-07-10T13:27:34</property>\n')
        block3.append('    <property name="xml">was here</property>\n')
        block3.append('    <property name="shotcut:caption">{slidename}</property>\n'.format(slidename=slidename))
        block3.append('  </producer>\n')
        #
        block4.append('    <entry producer="producer{id}" in="00:00:00.000" out="{slideduration}"/>\n'.format(id=id+len(files),slideduration=slideduration))
        #
        id=id+1
    #
    block2.append('  </playlist>\n')
    block2.append('  <producer id="black" in="00:00:00.000" out="{duration}">\n'.format(duration=duration))
    block2.append('    <property name="length">{duration}</property>\n'.format(duration=duration))
    block2.append('    <property name="eof">pause</property>\n')
    block2.append('    <property name="resource">0</property>\n')
    block2.append('    <property name="aspect_ratio">1</property>\n')
    block2.append('    <property name="mlt_service">color</property>\n')
    block2.append('    <property name="mlt_image_format">rgba</property>\n')
    block2.append('    <property name="set.test_audio">0</property>\n')
    block2.append('  </producer>\n')
    block2.append('  <playlist id="background">\n')
    block2.append('    <entry producer="black" in="00:00:00.000" out="{duration}"/>\n'.format(duration=duration))
    block2.append('  </playlist>\n')
    #
    block4.append('  </playlist>\n')
    block4.append('  <chain id="chain1" out="{duration}">\n'.format(duration=duration))
    block4.append('    <property name="length">{duration}</property>\n'.format(duration=duration))
    block4.append('    <property name="eof">pause</property>\n')
    block4.append('    <property name="resource">{wavefilename}</property>\n'.format(wavefilename=wavefilename))
    block4.append('    <property name="mlt_service">avformat-novalidate</property>\n')
    block4.append('    <property name="seekable">1</property>\n')
    block4.append('    <property name="audio_index">0</property>\n')
    block4.append('    <property name="video_index">-1</property>\n')
    block4.append('    <property name="mute_on_pause">0</property>\n')
    block4.append('    <property name="shotcut:hash">{audioHash}</property>\n'.format(audioHash=audioHash))
    block4.append('    <property name="shotcut:caption">{wavefilename}</property>\n'.format(wavefilename=os.path.basename(wavefilename)))
    block4.append('  </chain>\n')
    block4.append('  <playlist id="playlist1">\n')
    block4.append('    <property name="shotcut:audio">1</property>\n')
    block4.append('    <property name="shotcut:name">A1</property>\n')
    block4.append('    <entry producer="chain1" in="00:00:00.000" out="{duration}"/>\n'.format(duration=duration))
    block4.append('  </playlist>\n')
    block4.append('  <tractor id="tractor0" title="Shotcut version 21.10.31" in="00:00:00.000" out="{duration}">\n'.format(duration=duration))
    block4.append('    <property name="shotcut">1</property>\n')
    block4.append('    <property name="shotcut:projectAudioChannels">2</property>\n')
    block4.append('    <property name="shotcut:projectFolder">0</property>\n')
    block4.append('    <track producer="background"/>\n')
    block4.append('    <track producer="playlist0"/>\n')
    block4.append('    <track producer="playlist1" hide="video"/>\n')
    block4.append('    <transition id="transition0">\n')
    block4.append('      <property name="a_track">0</property>\n')
    block4.append('      <property name="b_track">1</property>\n')
    block4.append('      <property name="mlt_service">mix</property>\n')
    block4.append('      <property name="always_active">1</property>\n')
    block4.append('      <property name="sum">1</property>\n')
    block4.append('    </transition>\n')
    block4.append('    <transition id="transition1">\n')
    block4.append('      <property name="a_track">0</property>\n')
    block4.append('      <property name="b_track">1</property>\n')
    block4.append('      <property name="version">0.9</property>\n')
    block4.append('      <property name="mlt_service">frei0r.cairoblend</property>\n')
    block4.append('      <property name="threads">0</property>\n')
    block4.append('      <property name="disable">1</property>\n')
    block4.append('    </transition>\n')
    block4.append('    <transition id="transition2">\n')
    block4.append('      <property name="a_track">0</property>\n')
    block4.append('      <property name="b_track">2</property>\n')
    block4.append('      <property name="mlt_service">mix</property>\n')
    block4.append('      <property name="always_active">1</property>\n')
    block4.append('      <property name="sum">1</property>\n')
    block4.append('    </transition>\n')
    block4.append('  </tractor>\n')
    block4.append('</mlt>\n')

    dir_path = os.path.dirname(os.path.realpath(__file__)).replace("\\","/")
    filename=short_name + str(song_number).zfill(3) + ".mlt"
    f = open(filename, "w")

    f.write('<?xml version="1.0" standalone="no"?>\n')
    f.write('<mlt LC_NUMERIC="C" version="7.1.0" title="Shotcut version 21.10.31" producer="main_bin">\n')
    f.write('  <profile description="PAL 4:3 DV or DVD" width="1920" height="1080" progressive="1" sample_aspect_num="1" sample_aspect_den="1" display_aspect_num="16" display_aspect_den="9" frame_rate_num="25" frame_rate_den="1" colorspace="709"/>\n')
    f.write('  <chain id="chain0" out="{duration}">\n'.format(duration=duration))
    f.write('    <property name="length">{duration}</property>\n'.format(duration=duration))
    f.write('    <property name="eof">pause</property>\n')
    f.write('    <property name="resource">{wavefilename}</property>\n'.format(wavefilename=wavefilename))
    f.write('    <property name="mlt_service">avformat-novalidate</property>\n')
    f.write('    <property name="seekable">1</property>\n')
    f.write('    <property name="audio_index">0</property>\n')
    f.write('    <property name="video_index">-1</property>\n')
    f.write('    <property name="mute_on_pause">0</property>\n')
    f.write('    <property name="xml">was here</property>\n')
    f.write('    <property name="shotcut:hash">{audioHash}</property>\n'.format(audioHash=audioHash))
    f.write('  </chain>\n')

    for block in block1:
        f.write(block)
    for block in block2:
        f.write(block)
    for block in block3:
        f.write(block)
    for block in block4:
        f.write(block)
    f.close();

    subprocess.Popen(["C:/Program Files/Shotcut/shotcut.exe", filename])

def getCredentials():
    credentials = None

    if os.path.exists("../automate-youtube/token.pickle"):
        with open("../automate-youtube/token.pickle", "rb") as token:
            credentials = pickle.load(token)

    if not credentials or not credentials.valid:
        if credentials and credentials.expired and credentials.refresh_token:
            credentials.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(
                "../automate-youtube/client_secrets.json", scopes=[
                    "https://www.googleapis.com/auth/youtube.readonly",
                    "https://www.googleapis.com/auth/youtube",
                    "https://www.googleapis.com/auth/youtube.force-ssl",
                    "https://www.googleapis.com/auth/youtubepartner",
                    "https://www.googleapis.com/auth/youtubepartner-channel-audit",
                    "https://www.googleapis.com/auth/youtube.upload",
                    "https://www.googleapis.com/auth/youtube.channel-memberships.creator",
                    "https://www.googleapis.com/auth/youtube.third-party-link.creator",
                    "https://www.googleapis.com/auth/youtube.download"
                ]
            )

            flow.run_local_server(
                port=8080,
                prompt='consent',
                authorization_prompt_message=''
            )

            credentials = flow.credentials
        with open("../automate-youtube/token.pickle", "wb") as token:
            pickle.dump(credentials, token)
        # print(credentials.to_json())
    return credentials

def get_lines(filename):
  file = open(filename, 'r', encoding="UTF-8")
  Lines = file.readlines()
  file.close()
  return Lines

projects = {
  "GS": {"title": "Govinda Shatakam", "playlist": "PLdBTiOEoG70n-z67yuw2yPFtqPIZf4kfG" },
  "RGS": {"title": "Raghava Shatakam", "playlist": "PLdBTiOEoG70mYgAGZiaXxfctNofGMElJt" },
  "RS": {"title": "Radhika Shatakam", "playlist": "PLdBTiOEoG70n_aTSZgIGi3pOtyO7D2Tky" },
  "GE": {"title": "Geethavali", "playlist": "PLdBTiOEoG70mGzfaL26h4aFYV3eT0mWpb" },
  "KV": {"title": "Keerthanavali", "playlist": "" },
  "TS": {"title": "Teertha Shatakam", "playlist": "" },
  "YS": {"title": "Yugala Shatakam", "playlist": "" },
  "NGS": {"title": "Ranga Shatakam", "playlist": "" },
  "LSM": {"title": "Laghusthotramala", "playlist": "" },
  "RMS": {"title": "Radha Madhava Shatakam", "playlist": "" }
}

def upload_video(short_name='GS', long_name='govindashathakam'):
    song_number = get_song_number(long_name)
    ASTRING=short_name + str(song_number).zfill(3)
    video_file = ASTRING + ".mp4"
    thumbnail = ASTRING + "-Slide0001.png"

    if not path.exists(video_file):
        input("Video file does not exists. " + video_file)
        return

    index_file="C:/Users/rajar/books/srisrianna/"+long_name+"/index.txt"

    if not path.exists(index_file):
        input("Index file does not exists. " + index_file)
        return

    index_lines=get_lines(index_file)
    if len(index_lines) < song_number:
        input("Index file does not have title for song " + song_number)
        return

    title_meta=index_lines[song_number-1].rstrip("\n").split("|")
    raga_meta=title_meta[1].split(" - ")
    ragam=""
    talam=""
    if len(raga_meta) > 0:
        ragam=raga_meta[0]
    if len(raga_meta) > 1:
        talam=raga_meta[1]

    body = dict(
        snippet=dict(
            title="%s %s | %s | %s - %s" % (projects[short_name]['title'], str(song_number).zfill(3), title_meta[0], ragam, talam),
            description="""#%s #SriSriAnna #MythiliRajaraman

Song %s: %s
Ragam: %s 
Talam: %s
Composed by: Sri Sri Krishnapremi Swamigal

Playlist link: https://youtube.com/playlist?list=%s
            """ % (projects[short_name]['title'].replace(" ",""), str(song_number), title_meta[0], ragam, talam, projects[short_name]['playlist']),
            tags=[projects[short_name]['title'],
                  "premi",
                  "SriSriAnna",
                  "Shathakam",
                  "Satakam",
                  "bhajan",
                  "lyrics",
                  "krishnapremi",
                  "song " + str(song_number).zfill(3),
                  title_meta[0],
                  ragam
                 ],
            categoryId=10
        ),
        status=dict(
            privacyStatus="unlisted"
        )
    )

    with build('youtube', 'v3', static_discovery=False, credentials=getCredentials()) as service:
        try:
            # Call the API's videos.insert method to create and upload the video.
            insert_request = service.videos().insert(
                part=",".join(body.keys()),
                body=body,
                # The chunksize parameter specifies the size of each chunk of data, in
                # bytes, that will be uploaded at a time. Set a higher value for
                # reliable connections as fewer chunks lead to faster uploads. Set a lower
                # value for better recovery on less reliable connections.
                #
                # Setting "chunksize" equal to -1 in the code below means that the entire
                # file will be uploaded in a single HTTP request. (If the upload fails,
                # it will still be retried where it left off.) This is usually a best
                # practice, but if you're using Python older than 2.6 or if you're
                # running on App Engine, you should set the chunksize to something like
                # 1024 * 1024 (1 megabyte).
                media_body=MediaFileUpload(video_file, chunksize=-1, resumable=True)
            )

            uploaded_video_id = resumable_upload(insert_request)
            print('Setting thumbnail as ' + thumbnail)
            request = service.thumbnails().set(videoId=uploaded_video_id,media_body=MediaFileUpload(thumbnail))
            response = request.execute()
            print('Response for setting thumbnail: ')
            print(response)
            insert_playlistitem(projects[short_name]['playlist'], uploaded_video_id)
        except requests.HTTPError as e:
            print('Error response status code : {0}, reason : {1}'.format(
                e.status_code, e.error_details))

def resumable_upload(insert_request):
    # This method implements an exponential backoff strategy to resume a
    # failed upload.

    # api_key="AIzaSyDpvaVAspi9yIM8vVCPzajT42tdroYkULg"

    # Explicitly tell the underlying HTTP transport library not to retry, since
    # we are handling retry logic ourselves.
    httplib2.RETRIES = 1

    # Maximum number of times to retry before giving up.
    MAX_RETRIES = 10

    # Always retry when these exceptions are raised.
    RETRIABLE_EXCEPTIONS = (httplib2.HttpLib2Error, IOError, http.client.NotConnected,
                            http.client.IncompleteRead, http.client.ImproperConnectionState,
                            http.client.CannotSendRequest, http.client.CannotSendHeader,
                            http.client.ResponseNotReady, http.client.BadStatusLine)

    # Always retry when an apiclient.errors.HttpError with one of these status
    # codes is raised.
    RETRIABLE_STATUS_CODES = [500, 502, 503, 504]

    response = None
    error = None
    retry = 0
    while response is None:
        try:
            print("Uploading file...")
            status, response = insert_request.next_chunk()
            if response is not None:
                if 'id' in response:
                    print("Video id 'https://youtu.be/%s' was successfully uploaded." % response['id'])
                    return response['id']
                else:
                    print("The upload failed with an unexpected response: %s" % response)
        except requests.HTTPError as e:
            if e.resp.status in RETRIABLE_STATUS_CODES:
                error = "A retriable HTTP error %d occurred:\n%s" % (e.resp.status,
                                                                     e.content)
            else:
                raise
        except RETRIABLE_EXCEPTIONS as e:
            error = "A retriable error occurred: %s" % e

        if error is not None:
            print(error)
            retry += 1
            if retry > MAX_RETRIES:
                print("No longer attempting to retry.")

            max_sleep = 2 ** retry
            sleep_seconds = random.random() * max_sleep
            print("Sleeping %f seconds and then retrying..." % sleep_seconds)
            time.sleep(sleep_seconds)
    return None

def insert_playlistitem(playlistIdStr, videoIdStr):
    with build('youtube', 'v3', static_discovery=False, credentials=getCredentials()) as service:
        print('Adding to playlist')
        body=dict(
            snippet=dict(
                playlistId=playlistIdStr,
                resourceId=dict(
                    kind='youtube#video',
                    videoId=videoIdStr
                )
            )
        )
        print(body)
        insert_request = service.playlistItems().insert(
                            part=",".join(body.keys()),
                            body=body
                         )
        response = insert_request.execute()
        print('Response for adding to playlist: ')
        print(response)
