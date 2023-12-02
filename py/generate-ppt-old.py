import collections.abc
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re
import sys
import csv
import subprocess
import os
import glob
import time
import os.path
from os import path
from pathlib import Path
from pptx.enum.shapes import MSO_SHAPE

#Create multiple blue highlighted content slide
def create_slides(lines):
  bullets = [0.61,1.43,2.25,3.06,3.89,4.71,5.53,6.35]
  left = Inches(0.4) 
  width = height = Inches(0.28)
  slidecount=0
  next_space = ""
  for count in lines:
    if " " != count[0]:
      slide = prs.slides.add_slide(bullet_slide)
      slidecount = slidecount + 1
      shapes = slide.shapes
      #title_shape = shapes.title
      body_shape = shapes.placeholders[1]
      #title_shape.text = titletext
      tf = body_shape.text_frame
      p = tf.paragraphs[0]
      space = ""
      defa_color=RGBColor(0x6E, 0x6E, 0x6E)
      othr_color=RGBColor(0x6E, 0x6E, 0x6E)
      l=0
      for row in lines:
        if " " != row[0]:
          oval = shapes.add_shape(MSO_SHAPE.OVAL, left, Inches(bullets[l]), width, height)
        line_color = defa_color
        next_space = ""
        if not ("।" in row[0] or "॥" in row[0]):
          next_space = "    "
        if " " == row[0]:
          next_space = ""
        if row[0].strip().endswith("वाच") or row[0].strip().endswith("ऊचुः") or re.search(r"इति श्रीमद्भागवते महापुराणे .*पारमहंस्यां संहितायां", row[0]):
          line_color = othr_color
          next_space = ""
        if row == count:
          oval.fill.gradient()
          oval.fill.gradient_stops[0].color.rgb = RGBColor(0xFF, 0xFF, 0x00)
          oval.fill.gradient_stops[1].color.rgb = RGBColor(0xFF, 0xC0, 0x00)
          line_color = RGBColor(0xFF, 0xFF, 0x00)
          defa_color = RGBColor(0xFF, 0xFF, 0xFF)
          othr_color = RGBColor(0x0, 0xB0, 0xF0)
        run = p.add_run()
        run.text = space + row[0] + "\n"
        font = run.font
        font.color.rgb = line_color
        space = next_space
        l+=1
  return slidecount
#end create_slides

#
# Main Code Starts
#
canto="1"
if len(sys.argv) < 2:
  file_adyayam_number=""
  if path.exists("adyayam_number.txt"):
    f=open("adyayam_number.txt")
    file_adyayam_number = f.readlines()[0]
    f.close()
  adyayam_number = input ("Canto " + canto.zfill(2) + " Adhyayam Number?(" + file_adyayam_number + "): ")
  if adyayam_number == "" and file_adyayam_number != "":
    adyayam_number = file_adyayam_number
else:
  adyayam_number = sys.argv[1]
f=open("adyayam_number.txt", "w")
f.write(adyayam_number)
f.close()

titletext='Srimad Bhagavatham | Canto ' + canto.zfill(2) + ' | Chapter ' + adyayam_number.zfill(2)
prs = Presentation('template.pptm')
bullet_slide = prs.slide_layouts[1]
title_slide = prs.slide_layouts[0]
last_line=""

picture_file_name = 'C:/Users/rajar/Pictures/Bhagavatham' + canto.zfill(2) + '/' + adyayam_number.zfill(2) + '.png'
if not path.exists(picture_file_name):
  input("Picture file does not exists. " + picture_file_name)
  sys.exit()

ASTRING='SB-' + canto.zfill(2) + '-' + adyayam_number.zfill(2)
GENERATED_PPTM=ASTRING + ".pptm"

## First Slide
slide = prs.slides.add_slide(title_slide)
shapes = slide.shapes
picture = shapes.placeholders[10].insert_picture(picture_file_name)
shapes.placeholders[15].text_frame.text = "श्रीमद्भागवतमहापुराणम्"
shapes.placeholders[14].text_frame.text = "प्रथम स्कन्दः"
body_shape = shapes.placeholders[12] #Adyayam Text added later
shapes.placeholders[16].text_frame.text = canto.zfill(2)
shapes.placeholders[17].text_frame.text = adyayam_number.zfill(2)

#
# Debug only: For checking contents of the slide
#for shape in slide.placeholders:
#  print('%d %s' % (shape.placeholder_format.idx, shape.name))

## Main Loop to add content slides
lyrics_file='C:/Users/rajar/srimad_bhaghavatham/sanskrit/canto' + canto.zfill(2) + '/chapter'+adyayam_number.zfill(2)+".txt"
if not path.exists(lyrics_file):
  input("Lyrics file does not exists. " + lyrics_file)
  sys.exit()
lyricsdata = csv.reader(open(lyrics_file, encoding="UTF-8"))

firstrow=True
line=0
count=0
lines=[]
slidecount=1
for row in lyricsdata:
  if len(row) > 0:
    if firstrow:
      body_shape.text_frame.text = row[0]
      if len(sys.argv) == 4:
        prs.save(GENERATED_PPTM)
        p = subprocess.Popen(["C:/Program Files/Microsoft Office 15/root/office15/POWERPNT.EXE", "/M", GENERATED_PPTM, "Save_PowerPoint_Slide_as_Images"])
        sys.exit()
      firstrow=False
    else:
      if re.search(r"इति श्रीमद्भागवते महापुराणे .*पारमहंस्यां संहितायां", row[0]):
        last_line = row[0]
      else:
        lines.append(row)
        line = line + 1
        if line > 7:
          #Create multiple blue highlighted content slide
          slidecount = slidecount + create_slides(lines)
          line = 0
          lines=[]
if len(lines) > 0:
  slidecount = slidecount + create_slides(lines)

## Last Ending Slide
slide = prs.slides.add_slide(title_slide)
slidecount = slidecount + 1
shapes = slide.shapes
picture = shapes.placeholders[10].insert_picture(picture_file_name)

m=re.search(r"(इति श्रीमद्भागवते महापुराणे .*पारमहंस्यां संहितायां)\s(.*)\s(.*ऽध्यायः\s.*)", last_line)
if m == None:
  input("Last line of chapter" + adyayam_number.zfill(2) + ".txt is not in known format. Press any key to continue")
  sys.exit()
else:
  shapes.placeholders[13].text_frame.text = m.group(1)
  shapes.placeholders[14].text_frame.text = m.group(2)
  shapes.placeholders[12].text_frame.text = m.group(3)

prs.save(GENERATED_PPTM)

print ("Generated. Slide Count = ", slidecount)
if len(sys.argv) > 2:
  next_action=sys.argv[2]
else:
  next_action=input ("Export JPG (e) / Open (o): ")
if next_action == "e":
  last_slide = ((ASTRING + '-Slide%04d.jpg') % slidecount)
  files = glob.glob(ASTRING + '-Slide????.jpg')
  for f in files:
    try:
        os.remove(f)
    except OSError as e:
        print("Error: %s : %s" % (f, e.strerror))
  p = subprocess.Popen(["C:/Program Files/Microsoft Office 15/root/office15/POWERPNT.EXE", "/M", GENERATED_PPTM, "Save_PowerPoint_Slide_as_Images"])
  while not Path(last_slide).is_file():
    time.sleep(1)
  p.terminate()
  p.wait()
elif next_action == "o":
  p = subprocess.Popen(["C:/Program Files/Microsoft Office 15/root/office15/POWERPNT.EXE", GENERATED_PPTM])
